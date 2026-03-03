import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

# Read full content from the Linea de tiempo PPT
path = r"c:\Users\dvaldeb\Downloads\Linea de tiempo negociaciones- Comunicacion operacion1.pptx"
prs = Presentation(path)
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            texts.append(shape.text.strip())
        if shape.has_table:
            table = shape.table
            print(f'\n--- Slide {i+1} TABLE ({len(list(table.rows))} rows x {len(table.columns)} cols) ---', flush=True)
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                print(f'  Row {row_idx}: {cells}', flush=True)
    if texts:
        print(f'\n--- Slide {i+1} TEXT ---', flush=True)
        for t in texts:
            print(t, flush=True)
            print('---', flush=True)
