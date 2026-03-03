import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

files = [
    r'C:\Users\dvaldeb\Walmart Inc\ECO & LPM Desplegados - General\RRSS - RRLL\RELACIONES SINDICALES\Negociaciones Colectivas\Negociaciones 2025\FWM Chile Proceso Negociaci\u00f3n colectiva 2025\Informaci\u00f3n NC FWM Chile\Linea de tiempo negociaciones FW 1.pptx',
    r'C:\Users\dvaldeb\OneDrive - Walmart Inc\Respaldo notebook\ECO (1)\Relacionamiento Sindical\Material\Material 2026\Linea de tiempo negociaciones- Comunicacion operacion.pptx',
    r'C:\Users\dvaldeb\OneDrive - Walmart Inc\Respaldo notebook\ECO (1)\Relacionamiento Sindical\Material\Material 2026\Cierre Neg Colectivas 2025 y escenario 2026 (1).pptx',
    r'C:\Users\dvaldeb\OneDrive - Walmart Inc\Respaldo notebook\ECO (1)\Relacionamiento Sindical\Material\Material 2026\Equipos de emergencia (002) (1).pptx',
    r'C:\Users\dvaldeb\OneDrive - Walmart Inc\Respaldo notebook\ECO (1)\Relacionamiento Sindical\Material\Material 2026\Jornada LPT_21012026.pptx',
]

for path in files:
    print(f'\n{"="*60}', flush=True)
    print(f'FILE: {os.path.basename(path)}', flush=True)
    print('='*60, flush=True)
    if not os.path.exists(path):
        print('  NOT FOUND', flush=True)
        continue
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                print(f'\n--- Slide {i+1} ---', flush=True)
                for t in texts:
                    print(f'  {t[:300]}', flush=True)
    except Exception as e:
        print(f'  ERROR: {e}', flush=True)
