# -*- coding: utf-8 -*-
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import tempfile

pdf_path = r'c:\Users\dvaldeb\Downloads\Proceso Uniforme Modelo To- Be propuesta final 20251204 (1).pdf'
pptx_path = r'c:\Users\dvaldeb\pupy inteeligente\Diapositiva7_EDITABLE.pptx'

print('=== CREANDO DIAPOSITIVA 7 EDITABLE ===')
print()

# Abrir PDF
pdf = fitz.open(pdf_path)
page = pdf[6]  # Pagina 7

# Dimensiones
pdf_w = page.rect.width
pdf_h = page.rect.height
print(f'PDF: {pdf_w:.0f} x {pdf_h:.0f} pts')

# Crear PPT
prs = Presentation()

# Configurar slide (mismo aspect ratio que PDF)
slide_w = 13.333  # inches
slide_h = slide_w * (pdf_h / pdf_w)
prs.slide_width = Inches(slide_w)
prs.slide_height = Inches(slide_h)
print(f'PPT: {slide_w:.2f} x {slide_h:.2f} inches')

# Crear slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Funcion para convertir coordenadas
def to_ppt(x, y, w=None, h=None):
    px = (x / pdf_w) * slide_w
    py = (y / pdf_h) * slide_h
    if w and h:
        pw = (w / pdf_w) * slide_w
        ph = (h / pdf_h) * slide_h
        return px, py, pw, ph
    return px, py

# 1. AGREGAR IMAGEN DE FONDO (para mantener graficos perfectos)
print('Renderizando pagina como imagen de fondo...')
zoom = 3.0  # Alta resolucion
mat = fitz.Matrix(zoom, zoom)
pix = page.get_pixmap(matrix=mat)

temp_dir = tempfile.mkdtemp()
bg_path = os.path.join(temp_dir, 'page7_bg.png')
pix.save(bg_path)

# Agregar imagen de fondo
slide.shapes.add_picture(
    bg_path,
    Inches(0), Inches(0),
    Inches(slide_w), Inches(slide_h)
)
print('  Imagen de fondo agregada')

# 2. AGREGAR TEXTOS EDITABLES ENCIMA
print('Agregando textos editables...')

# Color principal del documento
DARK_BLUE = RGBColor(0, 30, 96)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(255, 0, 0)

# Titulo principal
x, y, w, h = to_ppt(45, 96, 280, 40)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Presupuesto (As-Is)"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
print('  + Titulo: Presupuesto (As-Is)')

# Descripcion
x, y, w, h = to_ppt(66, 140, 320, 50)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Descripcion:"
p.font.size = Pt(7)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p = tf.add_paragraph()
p.text = "El proceso de presupuesto para la compra de uniformes tiene como objetivo planificar y controlar los recursos financieros destinados a la adquisicion de uniformes para los colaboradores de la empresa, asegurando que dicha inversion se ajuste a las politicas internas."
p.font.size = Pt(7)
p.font.color.rgb = DARK_BLUE
print('  + Descripcion agregada')

# Dolores (Pain points)
x, y, w, h = to_ppt(413, 140, 260, 65)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Dolores:"
p.font.size = Pt(7)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

dolores = [
    "1. No existe registro de inventario.",
    "2. No tenemos registrado el detalle de las tallas de los asociados.",
    "3. No se incluye el dato de rotacion en la proyeccion.",
    "4. No se incluyen campanas en la proyeccion.",
    "5. No se incluyen inputs de negocio en la proyeccion.",
    "6. No existe control de presupuesto vs gasto real de uniformes."
]

for dolor in dolores:
    p = tf.add_paragraph()
    p.text = dolor
    p.font.size = Pt(7)
    p.font.color.rgb = DARK_BLUE
print('  + Dolores (6 items) agregados')

# Nota al pie
x, y, w, h = to_ppt(69, 459, 400, 15)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Nota: Etapa en constante cambio, no contamos con un proceso estandarizado."
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = DARK_BLUE
print('  + Nota al pie agregada')

# Labels Si/No en el diagrama
x, y = to_ppt(512, 386)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.3), Inches(0.2))
p = tb.text_frame.paragraphs[0]
p.text = "Si"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

x, y = to_ppt(459, 419)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.3), Inches(0.2))
p = tb.text_frame.paragraphs[0]
p.text = "No"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = RED
print('  + Labels Si/No agregados')

# Limpiar temporal
os.remove(bg_path)
os.rmdir(temp_dir)

# Guardar
print()
print(f'Guardando: {pptx_path}')
prs.save(pptx_path)

pdf.close()

print()
print('=== DIAPOSITIVA 7 CREADA ===')
print('Contenido:')
print('  - Imagen de fondo: Diagrama de flujo completo')
print('  - Textos editables: Titulo, Descripcion, Dolores, Nota, Labels')
print()
print('Para editar: Haz clic en cualquier texto para modificarlo.')
print('El diagrama de flujo se mantiene intacto como imagen de fondo.')
