import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt

files = [
    r"c:\Users\dvaldeb\Downloads\Linea de tiempo negociaciones- Comunicacion operacion1.pptx",
    r"c:\Users\dvaldeb\Downloads\Equipos de emergencia (002)1.pptx",
    r"c:\Users\dvaldeb\Downloads\Cierre Neg Colectivas 2025 y escenario 2026 1.pptx",
]

for path in files:
    print(f'\n{"="*80}', flush=True)
    print(f'FILE: {os.path.basename(path)}', flush=True)
    print(f'EXISTS: {os.path.exists(path)}', flush=True)
    print('='*80, flush=True)
    if not os.path.exists(path):
        continue
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
                # Check for tables
                if shape.has_table:
                    table = shape.table
                    print(f'\n--- Slide {i+1} TABLE ({table.rows.__len__()} rows x {len(table.columns)} cols) ---', flush=True)
                    for row_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        print(f'  Row {row_idx}: {cells}', flush=True)
            if texts:
                print(f'\n--- Slide {i+1} TEXT ---', flush=True)
                for t in texts:
                    print(f'  {t[:500]}', flush=True)
    except Exception as e:
        print(f'  ERROR: {e}', flush=True)
