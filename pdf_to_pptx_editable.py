# -*- coding: utf-8 -*-
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import tempfile
import io

# Rutas
pdf_path = r'c:\Users\dvaldeb\Downloads\Proceso Uniforme Modelo To- Be propuesta final 20251204 (1).pdf'
pptx_path = r'c:\Users\dvaldeb\pupy inteeligente\Proceso_Uniforme_EDITABLE.pptx'

print('=== CONVIRTIENDO PDF A POWERPOINT EDITABLE ===')
print()

# Abrir PDF
print(f'Abriendo PDF: {os.path.basename(pdf_path)}')
pdf = fitz.open(pdf_path)
num_pages = len(pdf)
print(f'Paginas encontradas: {num_pages}')
print()

# Obtener dimensiones de la primera pagina para referencia
first_page = pdf[0]
pdf_width = first_page.rect.width
pdf_height = first_page.rect.height
print(f'Dimensiones PDF: {pdf_width:.0f} x {pdf_height:.0f} pts')

# Crear PowerPoint
print('Creando PowerPoint...')
prs = Presentation()

# Configurar tamano del slide para coincidir con proporcion del PDF
# Convertir puntos PDF a pulgadas (72 pts = 1 inch)
slide_width_inches = pdf_width / 72
slide_height_inches = pdf_height / 72

# Escalar a tamano razonable (max 13.333 x 7.5 para 16:9)
scale = min(13.333 / slide_width_inches, 7.5 / slide_height_inches)
slide_width_inches *= scale
slide_height_inches *= scale

prs.slide_width = Inches(slide_width_inches)
prs.slide_height = Inches(slide_height_inches)

print(f'Dimensiones PPT: {slide_width_inches:.2f} x {slide_height_inches:.2f} inches')
print()

# Carpeta temporal para imagenes
temp_dir = tempfile.mkdtemp()

# Funcion para convertir coordenadas PDF a PPT
def pdf_to_ppt_coords(x, y, width, height):
    # PDF usa origen arriba-izquierda, igual que PPT
    # Escalar coordenadas
    ppt_x = (x / pdf_width) * slide_width_inches
    ppt_y = (y / pdf_height) * slide_height_inches
    ppt_w = (width / pdf_width) * slide_width_inches
    ppt_h = (height / pdf_height) * slide_height_inches
    return ppt_x, ppt_y, ppt_w, ppt_h

# Procesar cada pagina
print('Extrayendo elementos editables...')
for page_num, page in enumerate(pdf):
    print(f'  Pagina {page_num + 1}/{num_pages}:', end=' ')
    
    # Crear diapositiva en blanco
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    elements_count = 0
    
    # 1. Extraer imagenes de la pagina
    image_list = page.get_images(full=True)
    for img_index, img in enumerate(image_list):
        try:
            xref = img[0]
            base_image = pdf.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            
            # Guardar imagen temporal
            img_path = os.path.join(temp_dir, f'img_{page_num}_{img_index}.{image_ext}')
            with open(img_path, 'wb') as f:
                f.write(image_bytes)
            
            # Obtener posicion de la imagen en la pagina
            # Buscar en los contenidos de la pagina
            img_rects = page.get_image_rects(xref)
            if img_rects:
                rect = img_rects[0]
                x, y, w, h = pdf_to_ppt_coords(rect.x0, rect.y0, rect.width, rect.height)
                
                # Agregar imagen al slide
                slide.shapes.add_picture(
                    img_path,
                    Inches(x), Inches(y),
                    Inches(w), Inches(h)
                )
                elements_count += 1
        except Exception as e:
            pass  # Ignorar imagenes problematicas
    
    # 2. Extraer bloques de texto con posiciones
    text_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
    
    for block in text_dict.get("blocks", []):
        if block.get("type") == 0:  # Bloque de texto
            bbox = block.get("bbox", (0, 0, 0, 0))
            x0, y0, x1, y1 = bbox
            
            # Convertir coordenadas
            ppt_x, ppt_y, ppt_w, ppt_h = pdf_to_ppt_coords(x0, y0, x1 - x0, y1 - y0)
            
            # Minimo tamano
            if ppt_w < 0.1 or ppt_h < 0.1:
                continue
            
            # Crear textbox
            try:
                textbox = slide.shapes.add_textbox(
                    Inches(ppt_x), Inches(ppt_y),
                    Inches(max(ppt_w, 0.5)), Inches(max(ppt_h, 0.2))
                )
                tf = textbox.text_frame
                tf.word_wrap = True
                
                # Procesar lineas del bloque
                first_para = True
                for line in block.get("lines", []):
                    line_text = ""
                    font_size = 12
                    font_color = (0, 0, 0)
                    is_bold = False
                    
                    for span in line.get("spans", []):
                        line_text += span.get("text", "")
                        font_size = span.get("size", 12)
                        
                        # Extraer color
                        color = span.get("color", 0)
                        if isinstance(color, int):
                            r = (color >> 16) & 0xFF
                            g = (color >> 8) & 0xFF
                            b = color & 0xFF
                            font_color = (r, g, b)
                        
                        # Detectar bold
                        flags = span.get("flags", 0)
                        is_bold = bool(flags & 2 ** 4)  # Bold flag
                    
                    if line_text.strip():
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        
                        p.text = line_text.strip()
                        p.font.size = Pt(min(font_size * scale, 48))  # Escalar y limitar
                        p.font.bold = is_bold
                        try:
                            p.font.color.rgb = RGBColor(*font_color)
                        except:
                            p.font.color.rgb = RGBColor(0, 0, 0)
                
                elements_count += 1
            except Exception as e:
                pass
    
    # 3. Extraer dibujos/lineas (shapes)
    drawings = page.get_drawings()
    for drawing in drawings[:50]:  # Limitar a 50 por pagina
        try:
            rect = drawing.get("rect")
            if rect:
                x, y, w, h = pdf_to_ppt_coords(rect.x0, rect.y0, rect.width, rect.height)
                
                # Solo dibujar si es suficientemente grande
                if w > 0.05 and h > 0.05:
                    # Obtener color de relleno
                    fill_color = drawing.get("fill")
                    stroke_color = drawing.get("color")
                    
                    # Crear forma rectangulo
                    shape = slide.shapes.add_shape(
                        1,  # Rectangulo
                        Inches(x), Inches(y),
                        Inches(w), Inches(h)
                    )
                    
                    # Aplicar colores
                    if fill_color:
                        try:
                            r, g, b = int(fill_color[0]*255), int(fill_color[1]*255), int(fill_color[2]*255)
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(r, g, b)
                        except:
                            shape.fill.background()
                    else:
                        shape.fill.background()
                    
                    # Borde
                    if stroke_color:
                        try:
                            r, g, b = int(stroke_color[0]*255), int(stroke_color[1]*255), int(stroke_color[2]*255)
                            shape.line.color.rgb = RGBColor(r, g, b)
                        except:
                            pass
                    
                    elements_count += 1
        except:
            pass
    
    print(f'{elements_count} elementos')

# Limpiar archivos temporales
print()
print('Limpiando archivos temporales...')
for f in os.listdir(temp_dir):
    try:
        os.remove(os.path.join(temp_dir, f))
    except:
        pass
try:
    os.rmdir(temp_dir)
except:
    pass

pdf.close()

# Guardar PowerPoint
print(f'Guardando: {pptx_path}')
prs.save(pptx_path)

print()
print('=== CONVERSION COMPLETADA ===')
print(f'Archivo: {pptx_path}')
print(f'Diapositivas: {num_pages}')
print()
print('CONTENIDO EDITABLE:')
print('  - Textos: Puedes editar directamente')
print('  - Imagenes: Puedes mover, redimensionar, reemplazar')
print('  - Formas: Puedes modificar colores y tamanos')
