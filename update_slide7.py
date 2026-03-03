# -*- coding: utf-8 -*-
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import tempfile
import shutil

pdf_path = r'c:\Users\dvaldeb\Downloads\Proceso Uniforme Modelo To- Be propuesta final 20251204 (1).pdf'
pptx_src = r'c:\Users\dvaldeb\pupy inteeligente\Proceso_Uniforme_COMPLETO.pptx'
pptx_dst = r'c:\Users\dvaldeb\pupy inteeligente\Proceso_Uniforme_FINAL.pptx'

print('=== ACTUALIZANDO PPT CON DIAPOSITIVA 7 EDITABLE ===')
print()

# Copiar PPT original
shutil.copy2(pptx_src, pptx_dst)

# Abrir PDF para pagina 7
pdf = fitz.open(pdf_path)
page = pdf[6]

# Cargar PPT
prs = Presentation(pptx_dst)

# Obtener slide 7 (indice 6)
slide = prs.slides[6]

# Dimensiones
slide_w = prs.slide_width.inches
slide_h = prs.slide_height.inches
pdf_w = page.rect.width
pdf_h = page.rect.height

print(f'Slide 7: {slide_w:.2f} x {slide_h:.2f} inches')

def to_ppt(x, y, w=None, h=None):
    px = (x / pdf_w) * slide_w
    py = (y / pdf_h) * slide_h
    if w and h:
        pw = (w / pdf_w) * slide_w
        ph = (h / pdf_h) * slide_h
        return px, py, pw, ph
    return px, py

# Colores
DARK_BLUE = RGBColor(0, 30, 96)
RED = RGBColor(255, 0, 0)

print('Agregando textos editables sobre la imagen existente...')

# Titulo
x, y, w, h = to_ppt(45, 96, 280, 40)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
p = tb.text_frame.paragraphs[0]
p.text = "Presupuesto (As-Is)"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
print('  + Titulo')

# Descripcion
x, y, w, h = to_ppt(66, 140, 320, 50)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Descripcion:"
p.font.size = Pt(6)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p = tf.add_paragraph()
p.text = "El proceso de presupuesto para la compra de uniformes tiene como objetivo planificar y controlar los recursos financieros destinados a la adquisicion de uniformes para los colaboradores de la empresa."
p.font.size = Pt(6)
p.font.color.rgb = DARK_BLUE
print('  + Descripcion')

# Dolores
x, y, w, h = to_ppt(413, 140, 260, 65)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Dolores:"
p.font.size = Pt(6)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

for dolor in [
    "1. No existe registro de inventario.",
    "2. No tenemos registrado el detalle de las tallas.",
    "3. No se incluye el dato de rotacion.",
    "4. No se incluyen campanas en la proyeccion.",
    "5. No se incluyen inputs de negocio.",
    "6. No existe control de presupuesto vs gasto real."
]:
    p = tf.add_paragraph()
    p.text = dolor
    p.font.size = Pt(6)
    p.font.color.rgb = DARK_BLUE
print('  + Dolores (6)')

# Nota
x, y, w, h = to_ppt(69, 459, 400, 15)
tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
p = tb.text_frame.paragraphs[0]
p.text = "Nota: Etapa en constante cambio, no contamos con un proceso estandarizado."
p.font.size = Pt(7)
p.font.italic = True
p.font.color.rgb = DARK_BLUE
print('  + Nota')

# Labels Si/No
for label, pos, color in [("Si", (512, 386), DARK_BLUE), ("No", (459, 419), RED)]:
    x, y = to_ppt(*pos)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.3), Inches(0.2))
    p = tb.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = color
print('  + Labels Si/No')

pdf.close()

# Guardar
print()
print(f'Guardando: {pptx_dst}')
prs.save(pptx_dst)

print()
print('=== PPT ACTUALIZADO ===')
print('La diapositiva 7 ahora tiene:')
print('  - Imagen de fondo con el diagrama completo')
print('  - Textos editables superpuestos')
