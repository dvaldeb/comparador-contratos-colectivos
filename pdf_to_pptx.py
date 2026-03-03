# -*- coding: utf-8 -*-
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Rutas
pdf_path = r'c:\Users\dvaldeb\Downloads\Proceso Uniforme Modelo To- Be propuesta final 20251204 (1).pdf'
pptx_path = r'c:\Users\dvaldeb\pupy inteeligente\Proceso_Uniforme_To_Be.pptx'

print('=== CONVIRTIENDO PDF A POWERPOINT ===')
print()

# Leer PDF
print(f'Leyendo PDF: {os.path.basename(pdf_path)}')
reader = PdfReader(pdf_path)
num_pages = len(reader.pages)
print(f'Páginas encontradas: {num_pages}')
print()

# Extraer texto de cada página
pages_content = []
for i, page in enumerate(reader.pages):
    text = page.extract_text() or ''
    # Limpiar caracteres problemáticos
    text = text.encode('utf-8', errors='ignore').decode('utf-8')
    text = ''.join(c if ord(c) < 65536 else ' ' for c in text)
    pages_content.append(text)
    print(f'Página {i+1}: {len(text)} caracteres')

# Crear PowerPoint
print('Creando PowerPoint...')
prs = Presentation()

# Configurar tamaño de diapositiva (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colores Walmart
WALMART_BLUE = RGBColor(0, 83, 226)  # #0053e2
WALMART_SPARK = RGBColor(255, 194, 32)  # #ffc220

# Crear diapositiva de título
slide_layout = prs.slide_layouts[6]  # Layout en blanco
slide = prs.slides.add_slide(slide_layout)

# Título principal
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.5), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Proceso Uniforme"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WALMART_BLUE
p.alignment = PP_ALIGN.CENTER

# Subtítulo
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.5), Inches(1))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "Modelo To-Be - Propuesta Final"
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(100, 100, 100)
p.alignment = PP_ALIGN.CENTER

print('  Creada diapositiva de título')

# Crear diapositivas para cada página del PDF
for i, content in enumerate(pages_content):
    if not content or len(content.strip()) < 10:
        print(f'  Página {i+1}: Sin contenido significativo, saltando...')
        continue
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco
    
    # Título de la diapositiva
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    
    # Extraer primera línea como título
    lines = content.strip().split('\n')
    title_text = lines[0][:80] if lines else f'Página {i+1}'
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WALMART_BLUE
    
    # Contenido restante
    remaining_text = '\n'.join(lines[1:]) if len(lines) > 1 else content
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Dividir el contenido en párrafos manejables
    paragraphs = remaining_text.split('\n\n')
    
    for j, para_text in enumerate(paragraphs[:15]):  # Máximo 15 párrafos por slide
        if not para_text.strip():
            continue
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Limpiar y truncar texto si es muy largo
        clean_text = para_text.strip()[:500]
        p.text = clean_text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)
    
    print(f'  Creada diapositiva {i+2}')

# Guardar PowerPoint
print()
print(f'Guardando: {pptx_path}')
prs.save(pptx_path)

print()
print('=== CONVERSION COMPLETADA ===')
print(f'Archivo creado: {pptx_path}')
print(f'Total diapositivas: {len(prs.slides)}')
