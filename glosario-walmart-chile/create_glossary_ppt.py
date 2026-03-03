"""PPT v5 - Letra grande, resumenes ejecutivos, maximo impacto."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from glossary_ppt_data import (
    PRIMER_DIA, WALMART_GLOBAL, WALMART_CHILE, FORMATOS,
    MARCAS, SERVICIOS, KPIS, GLOSARIO,
)

# Colors
BLUE = RGBColor(0x00, 0x53, 0xE2)
BLUE_D = RGBColor(0x00, 0x33, 0x99)
BLUE_L = RGBColor(0xE8, 0xF0, 0xFF)
SPARK = RGBColor(0xFF, 0xC2, 0x20)
SPARK_D = RGBColor(0x99, 0x52, 0x13)
SPARK_L = RGBColor(0xFF, 0xF5, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GD = RGBColor(0x2D, 0x2D, 0x2D)
GM = RGBColor(0x74, 0x74, 0x74)
GL = RGBColor(0xF7, 0xF7, 0xF7)
GB = RGBColor(0xE0, 0xE0, 0xE0)
GRN = RGBColor(0x2A, 0x87, 0x03)
GRN_L = RGBColor(0xE8, 0xF7, 0xE0)
RED = RGBColor(0xEA, 0x11, 0x00)
PURPLE = RGBColor(0x6B, 0x21, 0xA8)
PURPLE_L = RGBColor(0xF3, 0xE8, 0xFF)
CYAN = RGBColor(0x06, 0x94, 0xA2)
CYAN_L = RGBColor(0xE0, 0xF7, 0xFA)
ORANGE = RGBColor(0xEA, 0x67, 0x00)
ORANGE_L = RGBColor(0xFF, 0xF0, 0xE0)
PINK = RGBColor(0xDB, 0x27, 0x77)
TEAL = RGBColor(0x0D, 0x96, 0x88)
TEAL_L = RGBColor(0xE0, 0xF2, 0xF1)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_L = RGBColor(0xEE, 0xF2, 0xFF)

COLOR_MAP = {
    "blue": (BLUE, BLUE_L), "purple": (PURPLE, PURPLE_L),
    "cyan": (CYAN, CYAN_L), "orange": (ORANGE, ORANGE_L),
    "indigo": (INDIGO, INDIGO_L), "teal": (TEAL, TEAL_L),
    "green": (GRN, GRN_L),
}

SW = Inches(13.333)
SH = Inches(7.5)


def _bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c


def _sh(s, tp, l, t, w, h, fill, border=None):
    sh = s.shapes.add_shape(tp, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    return sh


def _rect(s, l, t, w, h, fill, border=None):
    return _sh(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill, border)


def _box(s, l, t, w, h, fill):
    return _sh(s, MSO_SHAPE.RECTANGLE, l, t, w, h, fill)


def _circ(s, l, t, sz, fill):
    return _sh(s, MSO_SHAPE.OVAL, l, t, sz, sz, fill)


def _tx(s, l, t, w, h, text, sz=14, bold=False, color=GD, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = "Calibri"
    p.alignment = align
    return tb


def _ml(s, l, t, w, h, lines, sz=14, color=GD, sp=Pt(6), align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = sp if i > 0 else Pt(0); p.alignment = align


def _decor(s):
    _circ(s, Inches(-0.3), Inches(-0.3), Inches(1.0), SPARK)
    _circ(s, Inches(12.6), Inches(6.6), Inches(1.0), BLUE_L)


def _dots(s, x, y, c=SPARK_L):
    for r in range(3):
        for cc in range(4):
            _circ(s, x + cc * Inches(0.2), y + r * Inches(0.2), Inches(0.06), c)


def _hdr(s, title, accent=BLUE, sub=None):
    _box(s, Inches(0), Inches(0), SW, Inches(0.07), accent)
    _box(s, Inches(0), Inches(0.07), SW, Inches(0.035), SPARK)
    _circ(s, Inches(0.5), Inches(0.35), Inches(0.55), accent)
    _tx(s, Inches(0.5), Inches(0.37), Inches(0.55), Inches(0.5),
        "*", sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(1.2), Inches(0.3), Inches(11.4), Inches(0.6),
        title, sz=30, bold=True, color=GD)
    if sub:
        # Executive summary box
        _rect(s, Inches(0.5), Inches(0.95), Inches(12.2), Inches(0.65), SPARK_L, SPARK)
        _tx(s, Inches(0.7), Inches(1.0), Inches(11.8), Inches(0.55),
            sub, sz=16, bold=True, color=SPARK_D)


def _num(s, x, y, n, bg=BLUE, sz_circ=Inches(0.55)):
    _circ(s, x, y, sz_circ, bg)
    _tx(s, x, y + Inches(0.03), sz_circ, Inches(0.48),
        str(n), sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================

def s_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BLUE)
    _circ(s, Inches(-1.5), Inches(-1.5), Inches(5), BLUE_D)
    _circ(s, Inches(10), Inches(5), Inches(4), BLUE_D)
    _circ(s, Inches(5.5), Inches(0.3), Inches(0.8), SPARK)
    _circ(s, Inches(7.0), Inches(0.2), Inches(0.4), RGBColor(0xFF, 0xD7, 0x60))
    _tx(s, Inches(0), Inches(1.0), SW, Inches(0.9),
        "*", sz=70, bold=True, color=SPARK, align=PP_ALIGN.CENTER)
    _rect(s, Inches(2), Inches(2.2), Inches(9.3), Inches(2.6), RGBColor(0x00, 0x44, 0xBB))
    _tx(s, Inches(2.2), Inches(2.3), Inches(8.9), Inches(0.7),
        "Bienvenido/a a", sz=26, color=SPARK, align=PP_ALIGN.CENTER)
    _tx(s, Inches(2.2), Inches(3.0), Inches(8.9), Inches(1.0),
        "Walmart Chile", sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(2.2), Inches(3.9), Inches(8.9), Inches(0.6),
        "Tu guía de onboarding (inducción)", sz=18, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
    _rect(s, Inches(3.5), Inches(5.5), Inches(6.3), Inches(0.7), SPARK)
    _tx(s, Inches(3.5), Inches(5.55), Inches(6.3), Inches(0.6),
        "Nuevos colaboradores y alumnos en práctica | 2026",
        sz=14, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)


def s_bienvenida(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "Tu primer día empieza aquí", SPARK,
         PRIMER_DIA["bienvenida"])
    colors = [BLUE, GRN, PURPLE, ORANGE]
    tips = PRIMER_DIA["tips"]
    for i, (titulo, desc) in enumerate(tips):
        x = Inches(0.5) + i * Inches(3.1)
        y = Inches(2.0)
        _rect(s, x, y, Inches(2.9), Inches(2.8), WHITE, colors[i])
        _box(s, x, y, Inches(2.9), Inches(0.07), colors[i])
        _num(s, x + Inches(1.05), y + Inches(0.25), i + 1, colors[i], Inches(0.7))
        _tx(s, x + Inches(0.2), y + Inches(1.1), Inches(2.5), Inches(0.5),
            titulo, sz=16, bold=True, color=colors[i], align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.2), y + Inches(1.7), Inches(2.5), Inches(0.9),
            desc, sz=13, color=GD, align=PP_ALIGN.CENTER)

    # No te asustes
    _rect(s, Inches(0.5), Inches(5.2), Inches(12.2), Inches(2.0), SPARK_L, SPARK)
    _tx(s, Inches(0.7), Inches(5.25), Inches(11.8), Inches(0.5),
        "No te asustes si...", sz=18, bold=True, color=SPARK_D)
    items = PRIMER_DIA["no_asustes"]
    for i, item in enumerate(items):
        x = Inches(0.7) + i * Inches(3.05)
        _circ(s, x, Inches(5.85), Inches(0.35), colors[i])
        _tx(s, x, Inches(5.87), Inches(0.35), Inches(0.3),
            "!", sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.45), Inches(5.85), Inches(2.5), Inches(0.7),
            item, sz=12, color=SPARK_D)


def s_sabias_que(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "¿Sabías que...?", BLUE,
         "Datos curiosos para romper el hielo en tu primera semana.")
    facts = PRIMER_DIA["sabias_que"]
    colors = [BLUE, ORANGE, GRN, PURPLE]
    for i, fact in enumerate(facts):
        y = Inches(2.0) + i * Inches(1.3)
        c = colors[i]
        _rect(s, Inches(1.0), y, Inches(11.3), Inches(1.1), WHITE, c)
        _box(s, Inches(1.0), y, Inches(0.1), Inches(1.1), c)
        _circ(s, Inches(1.3), y + Inches(0.2), Inches(0.6), c)
        _tx(s, Inches(1.3), y + Inches(0.23), Inches(0.6), Inches(0.5),
            str(i + 1), sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, Inches(2.2), y + Inches(0.2), Inches(9.8), Inches(0.7),
            fact, sz=20, bold=True, color=c)


def s_walmart(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "¿Qué es Walmart?", BLUE,
         "La empresa minorista más grande del mundo. Fundada por Sam Walton.")
    # Metrics
    data = [
        ("1962", "Fundación", BLUE),
        ("~2.1M", "Colaboradores", GRN),
        ("#1", "Fortune 500", PURPLE),
        ("5,400+", "Tiendas Intl", ORANGE),
    ]
    for i, (val, lab, c) in enumerate(data):
        x = Inches(0.5) + i * Inches(3.1)
        _rect(s, x, Inches(2.0), Inches(2.9), Inches(1.5), c)
        _tx(s, x, Inches(2.05), Inches(2.9), Inches(0.8),
            val, sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, Inches(2.85), Inches(2.9), Inches(0.45),
            lab, sz=14, color=RGBColor(0xFF, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

    # Purpose
    _rect(s, Inches(0.5), Inches(3.8), Inches(12.2), Inches(1.2), SPARK)
    _circ(s, Inches(0.8), Inches(3.95), Inches(0.8), WHITE)
    _tx(s, Inches(0.8), Inches(4.0), Inches(0.8), Inches(0.7),
        "*", sz=30, bold=True, color=SPARK, align=PP_ALIGN.CENTER)
    _tx(s, Inches(1.8), Inches(3.9), Inches(10.6), Inches(0.45),
        "Nuestro Propósito:", sz=14, bold=True, color=BLUE_D)
    _tx(s, Inches(1.8), Inches(4.35), Inches(10.6), Inches(0.5),
        '"Ahorramos dinero a las personas para que vivan mejor."',
        sz=22, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)

    # Beliefs
    g = WALMART_GLOBAL
    bc = [SPARK, GRN, PURPLE]
    for i, (cr, desc) in enumerate(g["creencias"]):
        x = Inches(0.5) + i * Inches(4.15)
        y = Inches(5.3)
        _rect(s, x, y, Inches(3.95), Inches(1.6), WHITE, bc[i])
        _box(s, x, y, Inches(3.95), Inches(0.08), bc[i])
        _num(s, x + Inches(1.55), y + Inches(0.2), i + 1, bc[i], Inches(0.5))
        _tx(s, x + Inches(0.15), y + Inches(0.8), Inches(3.65), Inches(0.35),
            cr, sz=14, bold=True, color=bc[i], align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.15), y + Inches(1.15), Inches(3.65), Inches(0.35),
            desc, sz=11, color=GD, align=PP_ALIGN.CENTER)


def s_cultura(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "Cultura y Filosofía", SPARK,
         "Esto nos hace diferentes. Se vive en el día a día, no son solo palabras.")

    # EDLP + EDLC
    _rect(s, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.5), BLUE_L, BLUE)
    _circ(s, Inches(0.8), Inches(2.2), Inches(0.7), BLUE)
    _tx(s, Inches(0.8), Inches(2.25), Inches(0.7), Inches(0.6),
        "$", sz=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(1.7), Inches(2.2), Inches(4.5), Inches(0.5),
        "EDLP", sz=24, bold=True, color=BLUE)
    _tx(s, Inches(1.7), Inches(2.7), Inches(4.5), Inches(0.5),
        "Every Day Low Prices (Precios Bajos Todos los Días)", sz=14, color=BLUE_D)
    _tx(s, Inches(0.7), Inches(3.3), Inches(5.5), Inches(0.8),
        '"Piensa en cuán poco puedo cobrar, no en cuánto."\n- Sam Walton',
        sz=14, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)

    _rect(s, Inches(0.5), Inches(4.7), Inches(5.9), Inches(1.2), SPARK_L, SPARK)
    _circ(s, Inches(0.8), Inches(4.85), Inches(0.6), SPARK)
    _tx(s, Inches(0.8), Inches(4.88), Inches(0.6), Inches(0.5),
        "C", sz=22, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)
    _tx(s, Inches(1.6), Inches(4.85), Inches(4.6), Inches(0.4),
        "EDLC = Every Day Low Costs (Costos Bajos Todos los Días)", sz=16, bold=True, color=SPARK_D)
    _tx(s, Inches(1.6), Inches(5.3), Inches(4.6), Inches(0.5),
        "EDLC hace POSIBLE ofrecer EDLP.", sz=13, color=SPARK_D)

    # Culture practices
    practices = [
        ("Regla 3 Metros", "Para, sonríe, saluda.", BLUE),
        ("Regla Puesta de Sol", "Responde HOY mismo.", GRN),
        ("Puerta Abierta", "Sin miedo. Sin represalias.", PURPLE),
        ("Somos Colabora", "Socios, no empleados.", ORANGE),
    ]
    for i, (t, d, c) in enumerate(practices):
        y = Inches(2.0) + i * Inches(1.2)
        x = Inches(6.8)
        _rect(s, x, y, Inches(5.8), Inches(1.0), WHITE, c)
        _box(s, x, y, Inches(0.1), Inches(1.0), c)
        _circ(s, x + Inches(0.25), y + Inches(0.15), Inches(0.6), c)
        _tx(s, x + Inches(0.25), y + Inches(0.18), Inches(0.6), Inches(0.5),
            t[0], sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(1.0), y + Inches(0.1), Inches(4.5), Inches(0.4),
            t, sz=16, bold=True, color=c)
        _tx(s, x + Inches(1.0), y + Inches(0.55), Inches(4.5), Inches(0.4),
            d, sz=14, color=GD)

    # TIP
    _rect(s, Inches(6.8), Inches(6.2), Inches(5.8), Inches(0.7), GRN_L, GRN)
    _tx(s, Inches(7.0), Inches(6.25), Inches(5.4), Inches(0.6),
        "TIP: Se viven en tienda, oficina, tech, logística... en TODOS lados.",
        sz=13, bold=True, color=GRN)


def s_chile(prs):
    c = WALMART_CHILE
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "Walmart Chile", SPARK, c["historia"])

    metrics = [
        ("~400", "Tiendas", BLUE),
        ("~27%", "Mkt Share", GRN),
        ("#1", "Supermercado", PURPLE),
        ("+1,500", "Proveedores", ORANGE),
    ]
    for i, (val, lab, clr) in enumerate(metrics):
        x = Inches(0.5) + i * Inches(3.1)
        _rect(s, x, Inches(2.0), Inches(2.9), Inches(1.5), clr)
        _tx(s, x, Inches(2.05), Inches(2.9), Inches(0.8),
            val, sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, Inches(2.85), Inches(2.9), Inches(0.45),
            lab, sz=14, color=RGBColor(0xFF, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

    # Competencia
    _rect(s, Inches(0.5), Inches(3.8), Inches(5.9), Inches(3.3), GL, GB)
    _box(s, Inches(0.5), Inches(3.8), Inches(5.9), Inches(0.5), BLUE)
    _tx(s, Inches(0.7), Inches(3.83), Inches(5.5), Inches(0.4),
        "Competencia", sz=16, bold=True, color=WHITE)
    comp_c = [RED, ORANGE, CYAN, GRN]
    for i, (comp, d) in enumerate(c["competencia"]):
        y = Inches(4.5) + i * Inches(0.6)
        _circ(s, Inches(0.7), y + Inches(0.05), Inches(0.35), comp_c[i])
        _tx(s, Inches(1.2), y, Inches(2.0), Inches(0.4),
            comp, sz=14, bold=True, color=comp_c[i])
        _tx(s, Inches(3.2), y, Inches(3.0), Inches(0.4),
            d, sz=13, color=GD)

    # Sociedades SAP
    _rect(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.3), BLUE_L, BLUE)
    _tx(s, Inches(7.0), Inches(3.9), Inches(5.4), Inches(0.4),
        "Sociedades SAP", sz=16, bold=True, color=BLUE)
    _tx(s, Inches(7.0), Inches(4.3), Inches(5.4), Inches(0.3),
        "Los verás en todos los reportes:", sz=12, color=GM)
    for i, (cod, desc) in enumerate(c["sociedades"]):
        y = Inches(4.8) + i * Inches(0.45)
        _rect(s, Inches(7.0), y, Inches(5.4), Inches(0.38), WHITE)
        _tx(s, Inches(7.2), y + Inches(0.02), Inches(1.0), Inches(0.32),
            cod, sz=14, bold=True, color=BLUE)
        _tx(s, Inches(8.3), y + Inches(0.02), Inches(4.0), Inches(0.32),
            desc, sz=14, color=GD)


def s_estrategia(prs):
    c = WALMART_CHILE
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "7 Prioridades Estratégicas", BLUE,
         "Hacia dónde vamos como compañía. La visión a 5 años.")
    pr_c = [BLUE, GRN, PURPLE, ORANGE, CYAN, PINK, TEAL]
    for i, (t, d) in enumerate(c["prioridades"]):
        col = i // 4; row = i % 4
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(2.0) + row * Inches(1.25)
        c2 = pr_c[i]
        _rect(s, x, y, Inches(6.0), Inches(1.05), WHITE, c2)
        _box(s, x, y, Inches(0.1), Inches(1.05), c2)
        _num(s, x + Inches(0.25), y + Inches(0.18), i + 1, c2, Inches(0.6))
        _tx(s, x + Inches(1.05), y + Inches(0.1), Inches(4.7), Inches(0.45),
            t, sz=18, bold=True, color=c2)
        _tx(s, x + Inches(1.05), y + Inches(0.55), Inches(4.7), Inches(0.4),
            d, sz=14, color=GD)


def s_formatos(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "Formatos de Tienda", BLUE,
         "Cada formato atiende un tipo de cliente diferente. ¡Conócelos!")
    fc = [BLUE, GRN, PURPLE, ORANGE, GM]
    cw = Inches(2.35)
    for i, fmt in enumerate(FORMATOS):
        x = Inches(0.3) + i * (cw + Inches(0.15))
        y = Inches(2.0)
        c = fc[i]
        _rect(s, x, y, cw, Inches(4.8), GL, c)
        _rect(s, x, y, cw, Inches(0.7), c)
        _tx(s, x + Inches(0.1), y + Inches(0.08), cw - Inches(0.2), Inches(0.55),
            fmt["nombre"], sz=14, bold=True, color=WHITE)
        _rect(s, x + Inches(0.15), y + Inches(0.85), cw - Inches(0.3), Inches(0.4), c)
        _tx(s, x + Inches(0.15), y + Inches(0.87), cw - Inches(0.3), Inches(0.35),
            fmt["tipo"], sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.15), y + Inches(1.45), cw - Inches(0.3), Inches(0.8),
            fmt["desc"], sz=13, bold=True, color=GD, align=PP_ALIGN.CENTER)
        # NPS
        _rect(s, x + Inches(0.15), y + Inches(2.5), cw - Inches(0.3), Inches(0.7), SPARK_L)
        _tx(s, x + Inches(0.15), y + Inches(2.55), cw - Inches(0.3), Inches(0.25),
            "NPS Target", sz=9, bold=True, color=SPARK_D, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.15), y + Inches(2.8), cw - Inches(0.3), Inches(0.35),
            fmt["nps"], sz=16, bold=True, color=c, align=PP_ALIGN.CENTER)
        # SAP
        _rect(s, x + Inches(0.15), y + Inches(3.4), cw - Inches(0.3), Inches(0.5), BLUE_L)
        _tx(s, x + Inches(0.15), y + Inches(3.42), cw - Inches(0.3), Inches(0.2),
            "SAP", sz=9, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.15), y + Inches(3.62), cw - Inches(0.3), Inches(0.25),
            fmt["sap"], sz=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def s_marcas(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "Marcas Propias & Servicios", PURPLE,
         "Nuestras marcas ofrecen calidad a mejor precio. ¡Las reconocerás en góndola!")
    mc = [BLUE, GRN, PURPLE, ORANGE, CYAN, PINK]
    for i, (marca, cat) in enumerate(MARCAS):
        col = i // 3; row = i % 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(2.0) + row * Inches(1.15)
        c = mc[i]
        _rect(s, x, y, Inches(4.0), Inches(0.95), WHITE, c)
        _box(s, x, y, Inches(0.1), Inches(0.95), c)
        _circ(s, x + Inches(0.25), y + Inches(0.15), Inches(0.55), c)
        _tx(s, x + Inches(0.25), y + Inches(0.18), Inches(0.55), Inches(0.45),
            marca[0], sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.95), y + Inches(0.1), Inches(2.8), Inches(0.4),
            marca, sz=15, bold=True, color=c)
        _tx(s, x + Inches(0.95), y + Inches(0.52), Inches(2.8), Inches(0.35),
            cat, sz=12, color=GD)

    # Servicios
    _rect(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(3.45), GL, GB)
    _box(s, Inches(8.9), Inches(2.0), Inches(3.9), Inches(0.5), ORANGE)
    _tx(s, Inches(9.1), Inches(2.03), Inches(3.5), Inches(0.4),
        "Servicios Financieros", sz=14, bold=True, color=WHITE)
    sv_c = [RED, BLUE, SPARK_D, GRN]
    for i, (sv, desc) in enumerate(SERVICIOS):
        y = Inches(2.65) + i * Inches(0.65)
        _circ(s, Inches(9.1), y + Inches(0.05), Inches(0.35), sv_c[i])
        _tx(s, Inches(9.6), y, Inches(3.0), Inches(0.3),
            sv, sz=13, bold=True, color=sv_c[i])
        _tx(s, Inches(9.6), y + Inches(0.3), Inches(3.0), Inches(0.3),
            desc, sz=11, color=GD)

    # NPS bottom
    _rect(s, Inches(0.5), Inches(5.7), Inches(12.2), Inches(1.2), GRN_L, GRN)
    _circ(s, Inches(0.7), Inches(5.85), Inches(0.6), GRN)
    _tx(s, Inches(0.7), Inches(5.88), Inches(0.6), Inches(0.5),
        "NPS", sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(1.5), Inches(5.85), Inches(11.0), Inches(0.45),
        "NPS Targets (Metas) por Formato", sz=16, bold=True, color=GRN)
    _tx(s, Inches(1.5), Inches(6.3), Inches(11.0), Inches(0.45),
        "Líder >= 65%   |   Express >= 63%   |   SBA >= 78%   |   Mayorista >= 70%   |   Lider.cl >= 49%",
        sz=15, bold=True, color=GD)


def s_kpi(prs, sec, accent, light):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _dots(s, Inches(12.2), Inches(6.2), light)
    _hdr(s, f"KPI: {sec['titulo']}", accent, sec["resumen"])
    terms = sec["terms"]
    for i, (term, full, desc) in enumerate(terms):
        col = i // 2; row = i % 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(2.0) + row * Inches(2.55)
        _rect(s, x, y, Inches(6.0), Inches(2.3), light, accent)
        _box(s, x, y, Inches(0.1), Inches(2.3), accent)
        _circ(s, x + Inches(0.25), y + Inches(0.2), Inches(0.7), accent)
        _tx(s, x + Inches(0.25), y + Inches(0.23), Inches(0.7), Inches(0.6),
            term[0], sz=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(1.15), y + Inches(0.15), Inches(4.6), Inches(0.5),
            term, sz=22, bold=True, color=accent)
        _tx(s, x + Inches(1.15), y + Inches(0.65), Inches(4.6), Inches(0.35),
            full, sz=13, bold=True, color=GM)
        _tx(s, x + Inches(0.25), y + Inches(1.15), Inches(5.5), Inches(1.0),
            desc, sz=16, color=GD)


def s_gloss(prs, sec):
    accent, light = COLOR_MAP.get(sec["color"], (BLUE, BLUE_L))
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _dots(s, Inches(12.2), Inches(6.2), light)
    _hdr(s, sec["titulo"], accent, sec["resumen"])
    terms = sec["terms"]
    for i, (term, desc) in enumerate(terms):
        col = i // 2; row = i % 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(2.0) + row * Inches(2.55)
        _rect(s, x, y, Inches(6.0), Inches(2.3), light, accent)
        _box(s, x, y, Inches(0.1), Inches(2.3), accent)
        _circ(s, x + Inches(0.25), y + Inches(0.25), Inches(0.7), accent)
        _tx(s, x + Inches(0.25), y + Inches(0.28), Inches(0.7), Inches(0.6),
            term[0], sz=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(1.15), y + Inches(0.2), Inches(4.6), Inches(0.5),
            term, sz=22, bold=True, color=accent)
        _tx(s, x + Inches(0.25), y + Inches(0.95), Inches(5.5), Inches(1.2),
            desc, sz=16, color=GD)


def s_que_es_kpi(prs):
    """Slide explicativa: ¿Qué es un KPI?"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _decor(s)
    _hdr(s, "¿Qué es un KPI?", GRN,
         "Key Performance Indicator (Indicador Clave de Desempeño). Es un número que te dice si vas bien o mal.")

    # Analogía grande
    _rect(s, Inches(0.5), Inches(2.0), Inches(12.2), Inches(1.3), SPARK_L, SPARK)
    _circ(s, Inches(0.8), Inches(2.15), Inches(0.8), SPARK)
    _tx(s, Inches(0.8), Inches(2.2), Inches(0.8), Inches(0.7),
        "?", sz=32, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)
    _tx(s, Inches(1.8), Inches(2.1), Inches(10.6), Inches(0.45),
        "Piénsalo así:", sz=16, bold=True, color=SPARK_D)
    _tx(s, Inches(1.8), Inches(2.55), Inches(10.6), Inches(0.55),
        "Si tu auto tiene un tablero con velocímetro, bencina y temperatura... los KPIs son ESO para el negocio.",
        sz=18, bold=True, color=SPARK_D)

    # 4 cards explicativas
    cards = [
        ("¿Para qué sirven?", "Para saber si estamos cumpliendo\nlos objetivos del negocio.", BLUE),
        ("¿Quién los mira?", "TODOS. Desde el gerente de tienda\nhasta el CEO.", GRN),
        ("¿Cada cuánto?", "Diario, semanal o mensual\nsegún el indicador.", PURPLE),
        ("¿Qué hago con ellos?", "Si está en verde, bien. Si está en\nrojo, hay que actuar rápido.", ORANGE),
    ]
    for i, (titulo, desc, c) in enumerate(cards):
        x = Inches(0.5) + i * Inches(3.1)
        y = Inches(3.6)
        _rect(s, x, y, Inches(2.9), Inches(2.5), WHITE, c)
        _box(s, x, y, Inches(2.9), Inches(0.08), c)
        _num(s, x + Inches(1.05), y + Inches(0.2), i + 1, c, Inches(0.65))
        _tx(s, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(0.45),
            titulo, sz=15, bold=True, color=c, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.15), y + Inches(1.5), Inches(2.6), Inches(0.9),
            desc, sz=13, color=GD, align=PP_ALIGN.CENTER)

    # Semáforo
    _rect(s, Inches(0.5), Inches(6.4), Inches(12.2), Inches(0.8), GL, GB)
    _tx(s, Inches(0.7), Inches(6.45), Inches(2.5), Inches(0.35),
        "Semáforo de KPIs:", sz=14, bold=True, color=GD)
    _circ(s, Inches(3.5), Inches(6.5), Inches(0.5), GRN)
    _tx(s, Inches(4.1), Inches(6.52), Inches(2.0), Inches(0.35),
        "Verde = Cumple", sz=14, bold=True, color=GRN)
    _circ(s, Inches(6.3), Inches(6.5), Inches(0.5), SPARK)
    _tx(s, Inches(6.9), Inches(6.52), Inches(2.5), Inches(0.35),
        "Amarillo = Alerta", sz=14, bold=True, color=SPARK_D)
    _circ(s, Inches(9.5), Inches(6.5), Inches(0.5), RED)
    _tx(s, Inches(10.1), Inches(6.52), Inches(2.5), Inches(0.35),
        "Rojo = Actuar ya", sz=14, bold=True, color=RED)


def s_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BLUE)
    _circ(s, Inches(-1.5), Inches(-1.5), Inches(5), BLUE_D)
    _circ(s, Inches(10), Inches(5), Inches(4), BLUE_D)
    _circ(s, Inches(5.5), Inches(0.3), Inches(0.6), SPARK)
    _tx(s, Inches(0), Inches(0.8), SW, Inches(0.8),
        "*", sz=55, bold=True, color=SPARK, align=PP_ALIGN.CENTER)
    _rect(s, Inches(2), Inches(1.8), Inches(9.3), Inches(1.5), RGBColor(0x00, 0x44, 0xBB))
    _tx(s, Inches(2.2), Inches(1.85), Inches(8.9), Inches(0.7),
        "Ahora ya sabes lo básico.", sz=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(2.2), Inches(2.55), Inches(8.9), Inches(0.5),
        "El resto lo aprenderás haciendo. Y preguntando.",
        sz=18, color=SPARK, align=PP_ALIGN.CENTER)

    resources = [
        "MyWalmart - Portal interno",
        "Confluence - Glosarios y documentación",
        "Slack: #glosario_wt",
        "Workday - RRHH",
        "WIC+ - Dashboard KPIs",
        "Medallia - NPS y CSAT",
    ]
    _rect(s, Inches(2), Inches(3.5), Inches(9.3), Inches(2.4), RGBColor(0x00, 0x44, 0xBB))
    _tx(s, Inches(2.2), Inches(3.55), Inches(8.9), Inches(0.4),
        "Recursos útiles:", sz=16, bold=True, color=SPARK, align=PP_ALIGN.CENTER)
    _ml(s, Inches(2.5), Inches(4.0), Inches(8.3), Inches(1.8),
        resources, sz=14, color=RGBColor(0xCC, 0xDD, 0xFF),
        sp=Pt(5), align=PP_ALIGN.CENTER)

    _rect(s, Inches(3), Inches(6.2), Inches(7.3), Inches(0.65), SPARK)
    _tx(s, Inches(3), Inches(6.23), Inches(7.3), Inches(0.55),
        '"Solo hay un jefe... el cliente." - Sam Walton',
        sz=15, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)


# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Emu(int(SW))
    prs.slide_height = Emu(int(SH))

    s_cover(prs)
    s_bienvenida(prs)
    s_sabias_que(prs)
    s_walmart(prs)
    s_cultura(prs)
    s_chile(prs)
    s_estrategia(prs)
    s_formatos(prs)
    s_marcas(prs)

    for sec in GLOSARIO:
        s_gloss(prs, sec)

    s_que_es_kpi(prs)

    kpi_colors = [
        (PINK, RGBColor(0xFC, 0xE7, 0xF3)),
        (BLUE, BLUE_L),
        (GRN, GRN_L),
        (PURPLE, PURPLE_L),
        (RED, RGBColor(0xFF, 0xEB, 0xE8)),
        (CYAN, CYAN_L),
    ]
    for sec, (ac, lt) in zip(KPIS, kpi_colors):
        s_kpi(prs, sec, ac, lt)

    s_closing(prs)

    out = "Glosario_Walmart_Chile_v7.pptx"
    prs.save(out)
    print(f"Guardado: {out} | Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
