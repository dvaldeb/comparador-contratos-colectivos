# -*- coding: utf-8 -*-
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os
import tempfile

# Rutas
pdf_path = r'c:\Users\dvaldeb\Downloads\Proceso Uniforme Modelo To- Be propuesta final 20251204 (1).pdf'
pptx_path = r'c:\Users\dvaldeb\pupy inteeligente\Proceso_Uniforme_COMPLETO.pptx'

print('=== CONVIRTIENDO PDF COMPLETO A POWERPOINT ===')
print()

# Crear carpeta temporal para imagenes
temp_dir = tempfile.mkdtemp()
print(f'Carpeta temporal: {temp_dir}')
print()

# Abrir PDF
print(f'Abriendo PDF: {os.path.basename(pdf_path)}')
pdf = fitz.open(pdf_path)
num_pages = len(pdf)
print(f'Paginas encontradas: {num_pages}')
print()

# Crear PowerPoint
print('Creando PowerPoint...')
prs = Presentation()

# Configurar tamano 16:9 (igual que el PDF tipicamente)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Procesar cada pagina
print('Extrayendo paginas como imagenes...')
for i, page in enumerate(pdf):
    print(f'  Procesando pagina {i+1}/{num_pages}...', end=' ')
    
    # Renderizar pagina a imagen con alta resolucion
    # zoom=2 da buena calidad sin ser muy pesado
    zoom = 2.0
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat)
    
    # Guardar imagen temporal
    img_path = os.path.join(temp_dir, f'page_{i+1:03d}.png')
    pix.save(img_path)
    
    # Crear diapositiva
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Calcular dimensiones para ajustar la imagen
    # Mantener proporcion del PDF
    img_width = pix.width
    img_height = pix.height
    
    # Calcular escala para ajustar al slide
    slide_w = 13.333
    slide_h = 7.5
    
    # Calcular ratio
    ratio_w = slide_w / (img_width / 96)  # 96 DPI base
    ratio_h = slide_h / (img_height / 96)
    ratio = min(ratio_w, ratio_h) * 0.95  # 95% para margen
    
    final_w = (img_width / 96) * ratio
    final_h = (img_height / 96) * ratio
    
    # Centrar imagen
    left = (slide_w - final_w) / 2
    top = (slide_h - final_h) / 2
    
    # Agregar imagen al slide
    slide.shapes.add_picture(
        img_path,
        Inches(left),
        Inches(top),
        Inches(final_w),
        Inches(final_h)
    )
    
    print(f'OK ({pix.width}x{pix.height})')

# Guardar PowerPoint
print()
print(f'Guardando PowerPoint: {pptx_path}')
prs.save(pptx_path)

# Limpiar imagenes temporales
print('Limpiando archivos temporales...')
for f in os.listdir(temp_dir):
    os.remove(os.path.join(temp_dir, f))
os.rmdir(temp_dir)

pdf.close()

print()
print('=== CONVERSION COMPLETADA ===')
print(f'Archivo creado: {pptx_path}')
print(f'Total diapositivas: {num_pages}')
print()
print('El PowerPoint contiene cada pagina del PDF como imagen de alta calidad.')
print('Puedes editar agregando formas, texto, etc. sobre las imagenes.')
